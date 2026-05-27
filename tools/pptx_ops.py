import re
import subprocess
import sys
import tempfile
from pathlib import Path
from typing import Any, Optional

from pptx import Presentation
from pptx.util import Emu, Inches, Pt

from .file_ops import _validate_path
from .security import require_unlocked


# Layout indices in the default python-pptx blank template.
# 0=title slide, 1=title+content, 2=section header, 3=two content,
# 5=title only, 6=blank
LAYOUT_TITLE = 0
LAYOUT_TITLE_CONTENT = 1
LAYOUT_SECTION = 2
LAYOUT_TITLE_ONLY = 5
LAYOUT_BLANK = 6


def create_pptx(path: str, slides: list[dict], title: Optional[str] = None) -> str:
    """Create a .pptx file from a list of slide specs.

    Each slide is a dict supporting these keys (all optional unless noted):
      - title: str
      - subtitle: str            (only used for the cover slide)
      - bullets: list[str]       (one bullet per item; nested via "  " indent)
      - body: str                (multiline freeform text, used if no bullets)
      - image_path: str          (PNG/JPG path, placed on the right half)
      - notes: str               (presenter notes)
      - layout: "cover" | "title_content" | "section" | "title_only" | "blank"

    If layout is omitted, a sensible default is chosen.

    Args:
        path: Output .pptx path under the allowed base.
        slides: Ordered list of slide spec dicts. Must contain at least one slide.
        title: Optional deck title used for file metadata.
    """
    locked = require_unlocked()
    if locked:
        return locked
    try:
        if not isinstance(slides, list) or not slides:
            return "[create_pptx error: 'slides' must be a non-empty list]"
        out = _validate_path(path)
        if out.suffix.lower() != ".pptx":
            return "[create_pptx error: path must end with .pptx]"
        out.parent.mkdir(parents=True, exist_ok=True)

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        if title:
            prs.core_properties.title = title

        for spec in slides:
            if not isinstance(spec, dict):
                return "[create_pptx error: each slide must be an object]"
            _add_slide(prs, spec)

        prs.save(str(out))
        return f"Wrote PowerPoint {out} ({len(slides)} slide(s))"
    except Exception as e:
        return f"[create_pptx error: {type(e).__name__}: {e}]"


def pptx_from_markdown(path: str, markdown: str, title: Optional[str] = None) -> str:
    """Convert a markdown document into slides.

    Conventions:
      - The first H1 (# ...) becomes the cover slide title.
      - Each subsequent H1 starts a new slide (title + bullets).
      - H2 inside a slide becomes a sub-heading line.
      - Bullet lines (- or *) become slide bullets.
      - "---" creates a section divider slide using the next H1 as the title.

    Args:
        path: Output .pptx path under the allowed base.
        markdown: Markdown source text.
        title: Optional deck title used for file metadata.
    """
    slides = _markdown_to_slide_specs(markdown)
    if not slides:
        return "[pptx_from_markdown error: no headings found; nothing to render]"
    return create_pptx(path, slides, title=title or slides[0].get("title"))


def pptx_info(path: str) -> str:
    """Return a structural summary of an existing .pptx (slide count + titles)."""
    try:
        p = _validate_path(path)
        if p.suffix.lower() != ".pptx":
            return "[pptx_info error: path must end with .pptx]"
        prs = Presentation(str(p))
        lines = [f"slides: {len(prs.slides)}"]
        for i, slide in enumerate(prs.slides, 1):
            title_text = ""
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title_text = slide.shapes.title.text_frame.text.strip()
            shapes_with_text = sum(
                1 for shp in slide.shapes if shp.has_text_frame and shp.text_frame.text.strip()
            )
            lines.append(f"  {i:>3}. {title_text or '(untitled)'}  [{shapes_with_text} text shapes]")
        return "\n".join(lines)
    except Exception as e:
        return f"[pptx_info error: {type(e).__name__}: {e}]"


def pptx_add_slide(pptx_path: str, slide: dict) -> str:
    """Append one slide to an existing .pptx using the same spec shape as create_pptx.

    Args:
        pptx_path: Existing .pptx path.
        slide: Slide spec dict with keys title/subtitle/bullets/body/image_path/notes/layout.
    """
    locked = require_unlocked()
    if locked:
        return locked
    try:
        if not isinstance(slide, dict):
            return "[pptx_add_slide error: slide must be an object]"
        p = _validate_path(pptx_path)
        if not p.is_file():
            return f"[pptx_add_slide error: file not found: {p}]"
        prs = Presentation(str(p))
        _add_slide(prs, slide)
        prs.save(str(p))
        return f"Appended 1 slide to {p} (now {len(prs.slides)} total)"
    except Exception as e:
        return f"[pptx_add_slide error: {type(e).__name__}: {e}]"


def pptx_add_image(
    pptx_path: str,
    slide_index: int,
    image_path: str,
    left_in: float = 0.7,
    top_in: float = 1.4,
    width_in: Optional[float] = None,
    height_in: Optional[float] = None,
) -> str:
    """Place an image on an existing slide (1-indexed).

    Args:
        pptx_path: Existing .pptx path.
        slide_index: 1-based slide index to modify.
        image_path: Image file to insert.
        left_in: Position from left, in inches.
        top_in: Position from top, in inches.
        width_in: Width in inches. If both width and height are omitted, the
            image's native size is used.
        height_in: Height in inches.
    """
    locked = require_unlocked()
    if locked:
        return locked
    try:
        p = _validate_path(pptx_path)
        if not p.is_file():
            return f"[pptx_add_image error: file not found: {p}]"
        img = _validate_path(image_path)
        if not img.is_file():
            return f"[pptx_add_image error: image not found: {img}]"
        prs = Presentation(str(p))
        if slide_index < 1 or slide_index > len(prs.slides):
            return f"[pptx_add_image error: slide_index out of range (1..{len(prs.slides)})]"
        slide = prs.slides[slide_index - 1]
        kwargs = {}
        if width_in is not None:
            kwargs["width"] = Inches(width_in)
        if height_in is not None:
            kwargs["height"] = Inches(height_in)
        pic = slide.shapes.add_picture(str(img), Inches(left_in), Inches(top_in), **kwargs)
        prs.save(str(p))
        w_in = pic.width / 914400.0
        h_in = pic.height / 914400.0
        return (
            f"Added image on slide {slide_index} of {p}\n"
            f"  position: ({left_in:.2f}, {top_in:.2f}) in\n"
            f"  size: {w_in:.2f} x {h_in:.2f} in"
        )
    except Exception as e:
        return f"[pptx_add_image error: {type(e).__name__}: {e}]"


def pptx_replace_image(
    pptx_path: str,
    slide_index: int,
    new_image_path: str,
    image_index: int = 1,
) -> str:
    """Replace an existing picture on a slide with a new image, preserving position/size.

    Args:
        pptx_path: Existing .pptx path.
        slide_index: 1-based slide index.
        new_image_path: Replacement image file.
        image_index: 1-based index among picture shapes on the slide.
    """
    locked = require_unlocked()
    if locked:
        return locked
    try:
        p = _validate_path(pptx_path)
        if not p.is_file():
            return f"[pptx_replace_image error: file not found: {p}]"
        img = _validate_path(new_image_path)
        if not img.is_file():
            return f"[pptx_replace_image error: image not found: {img}]"
        prs = Presentation(str(p))
        if slide_index < 1 or slide_index > len(prs.slides):
            return f"[pptx_replace_image error: slide_index out of range (1..{len(prs.slides)})]"
        slide = prs.slides[slide_index - 1]
        pics = [shp for shp in slide.shapes if shp.shape_type == 13]  # 13 = PICTURE
        if not pics:
            return f"[pptx_replace_image error: no picture shapes on slide {slide_index}]"
        if image_index < 1 or image_index > len(pics):
            return f"[pptx_replace_image error: image_index out of range (1..{len(pics)})]"
        target = pics[image_index - 1]
        left, top, width, height = target.left, target.top, target.width, target.height
        sp = target._element
        sp.getparent().remove(sp)
        slide.shapes.add_picture(str(img), left, top, width=width, height=height)
        prs.save(str(p))
        return f"Replaced picture #{image_index} on slide {slide_index} of {p}"
    except Exception as e:
        return f"[pptx_replace_image error: {type(e).__name__}: {e}]"


def pptx_add_table(
    pptx_path: str,
    slide_index: int,
    rows: list[list[str]],
    left_in: float = 0.6,
    top_in: float = 1.4,
    width_in: float = 11.8,
    height_in: float = 5.0,
    header_row: bool = True,
) -> str:
    """Insert a table onto an existing slide.

    Args:
        pptx_path: Existing .pptx path.
        slide_index: 1-based slide index.
        rows: 2D list of cell strings. The first row is treated as header when header_row=True.
        left_in, top_in: Position in inches.
        width_in, height_in: Size in inches.
        header_row: When true, bold the first row.
    """
    locked = require_unlocked()
    if locked:
        return locked
    try:
        if not isinstance(rows, list) or not rows or not isinstance(rows[0], list):
            return "[pptx_add_table error: rows must be a 2D list]"
        p = _validate_path(pptx_path)
        if not p.is_file():
            return f"[pptx_add_table error: file not found: {p}]"
        prs = Presentation(str(p))
        if slide_index < 1 or slide_index > len(prs.slides):
            return f"[pptx_add_table error: slide_index out of range (1..{len(prs.slides)})]"
        slide = prs.slides[slide_index - 1]
        n_rows = len(rows)
        n_cols = max(len(r) for r in rows)
        tbl = slide.shapes.add_table(
            n_rows, n_cols, Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in)
        ).table
        for r, row in enumerate(rows):
            for c in range(n_cols):
                text = str(row[c]) if c < len(row) else ""
                cell = tbl.cell(r, c)
                cell.text = text
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(12)
                        if header_row and r == 0:
                            run.font.bold = True
        prs.save(str(p))
        return f"Added {n_rows}x{n_cols} table on slide {slide_index} of {p}"
    except Exception as e:
        return f"[pptx_add_table error: {type(e).__name__}: {e}]"


def pptx_export_png(
    pptx_path: str,
    output_dir: Optional[str] = None,
    width_px: int = 1280,
) -> str:
    """Render every slide of a .pptx to PNG using Microsoft PowerPoint (Windows COM).

    Use this to self-verify slide layout before reporting completion. Requires
    PowerPoint to be installed on the host.

    Args:
        pptx_path: Existing .pptx path.
        output_dir: Directory for the PNGs (created if missing). Defaults to
            "<pptx>_preview" alongside the pptx.
        width_px: Target PNG width in pixels.
    """
    try:
        p = _validate_path(pptx_path)
        if not p.is_file():
            return f"[pptx_export_png error: file not found: {p}]"
        out_dir = _validate_path(output_dir) if output_dir else p.with_suffix("").parent / f"{p.stem}_preview"
        if isinstance(out_dir, Path):
            out_dir = _validate_path(str(out_dir))
        out_dir.mkdir(parents=True, exist_ok=True)

        try:
            import pythoncom  # type: ignore
            import win32com.client  # type: ignore
        except ImportError:
            return "[pptx_export_png error: pywin32 not installed; cannot drive PowerPoint COM]"

        pythoncom.CoInitialize()
        ppt = None
        pres = None
        try:
            ppt = win32com.client.DispatchEx("PowerPoint.Application")
            # PowerPoint can't be fully hidden on all builds; the dispatch-export-quit pattern is the standard.
            pres = ppt.Presentations.Open(str(p), WithWindow=False)
            scale_h = int(width_px * (pres.PageSetup.SlideHeight / pres.PageSetup.SlideWidth))
            pres.Export(str(out_dir), "PNG", width_px, scale_h)
            count = len(list(out_dir.glob("*.PNG"))) + len(list(out_dir.glob("*.png")))
            return f"Exported {count} PNG(s) to {out_dir}"
        finally:
            try:
                if pres is not None:
                    pres.Close()
            except Exception:
                pass
            try:
                if ppt is not None:
                    ppt.Quit()
            except Exception:
                pass
            pythoncom.CoUninitialize()
    except Exception as e:
        return f"[pptx_export_png error: {type(e).__name__}: {e}]"


# ---------- internals ----------

def _add_slide(prs: Presentation, spec: dict) -> None:
    layout_name = spec.get("layout")
    layout_idx = _pick_layout(spec, layout_name)
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)

    title = spec.get("title")
    subtitle = spec.get("subtitle")
    bullets = spec.get("bullets")
    body = spec.get("body")
    image_path = spec.get("image_path")
    notes = spec.get("notes")

    if title and slide.shapes.title is not None:
        slide.shapes.title.text = str(title)
        _style_title(slide.shapes.title)

    if layout_idx == LAYOUT_TITLE and subtitle:
        placeholders = [ph for ph in slide.placeholders if ph.placeholder_format.idx == 1]
        if placeholders:
            placeholders[0].text = str(subtitle)

    # Find the content placeholder (idx 1) for text body if present.
    body_ph = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1 and ph != slide.shapes.title:
            body_ph = ph
            break

    if bullets and isinstance(bullets, list) and body_ph is not None:
        _fill_bullets(body_ph, bullets)
    elif body and body_ph is not None:
        _fill_body(body_ph, str(body))

    if image_path:
        img = _validate_path(image_path)
        if not img.is_file():
            raise FileNotFoundError(f"image_path not found: {img}")
        bullets_or_body = bullets or body
        if bullets_or_body:
            # Right-half layout: text on left, image on right.
            left = Inches(7.0)
            top = Inches(1.5)
            width = Inches(5.8)
            slide.shapes.add_picture(str(img), left, top, width=width)
        else:
            # Full-width hero image under the title.
            left = Inches(0.7)
            top = Inches(1.4)
            width = Inches(11.9)
            slide.shapes.add_picture(str(img), left, top, width=width)

    if notes:
        slide.notes_slide.notes_text_frame.text = str(notes)


def _pick_layout(spec: dict, explicit: Optional[str]) -> int:
    if explicit:
        return {
            "cover": LAYOUT_TITLE,
            "title_content": LAYOUT_TITLE_CONTENT,
            "section": LAYOUT_SECTION,
            "title_only": LAYOUT_TITLE_ONLY,
            "blank": LAYOUT_BLANK,
        }.get(explicit, LAYOUT_TITLE_CONTENT)
    if spec.get("subtitle"):
        return LAYOUT_TITLE
    if not spec.get("bullets") and not spec.get("body") and spec.get("title"):
        return LAYOUT_SECTION
    return LAYOUT_TITLE_CONTENT


def _style_title(title_shape: Any) -> None:
    tf = title_shape.text_frame
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(36)
            run.font.bold = True


def _fill_bullets(placeholder: Any, bullets: list) -> None:
    tf = placeholder.text_frame
    tf.clear()
    first = True
    for raw in bullets:
        text = str(raw)
        indent = 0
        # Honor a leading "  " (2-space) indent per level for nesting.
        while text.startswith("  "):
            text = text[2:]
            indent += 1
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.text = text
        para.level = min(indent, 4)
        for run in para.runs:
            run.font.size = Pt(20 - indent * 2)


def _fill_body(placeholder: Any, text: str) -> None:
    tf = placeholder.text_frame
    tf.clear()
    lines = text.splitlines() or [""]
    first = True
    for line in lines:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.text = line
        for run in para.runs:
            run.font.size = Pt(18)


def _markdown_to_slide_specs(md: str) -> list[dict]:
    lines = md.splitlines()
    slides: list[dict] = []
    current: Optional[dict] = None
    cover_done = False

    def flush() -> None:
        if current is not None:
            slides.append(current)

    section_pending = False

    for raw_line in lines:
        line = raw_line.rstrip()
        if line.strip() == "---":
            section_pending = True
            continue

        h1 = re.match(r"^#\s+(.*)", line)
        h2 = re.match(r"^##\s+(.*)", line)
        bullet = re.match(r"^\s*[-*]\s+(.*)", line)

        if h1:
            flush()
            heading = h1.group(1).strip()
            if not cover_done:
                current = {"title": heading, "layout": "cover"}
                cover_done = True
            elif section_pending:
                current = {"title": heading, "layout": "section"}
                section_pending = False
            else:
                current = {"title": heading, "bullets": [], "layout": "title_content"}
            continue

        if current is None:
            current = {"title": "Overview", "bullets": [], "layout": "title_content"}

        if h2:
            current.setdefault("bullets", []).append(h2.group(1).strip())
            continue

        if bullet:
            indent_spaces = len(raw_line) - len(raw_line.lstrip(" "))
            level = indent_spaces // 2
            current.setdefault("bullets", []).append("  " * level + bullet.group(1).strip())
            continue

        if line.strip():
            # Treat orphan paragraph lines as bullet items so they show up.
            current.setdefault("bullets", []).append(line.strip())

    flush()
    return slides
